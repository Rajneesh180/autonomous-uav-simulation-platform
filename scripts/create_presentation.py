#!/usr/bin/env python3
"""Generate a professional BTP Mid-Term Presentation (PPTX)."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "report" / "figures"
OUT = ROOT / "report" / "BTP_MidTerm_Presentation.pptx"

# ── Colour palette ──────────────────────────────────────────────────────
DARK = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy
ACCENT = RGBColor(0x00, 0x7A, 0xCC)  # blue accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x00, 0x96, 0x4B)


def _add_bg(slide, color=DARK):
    """Solid background fill."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_accent_bar(slide, top=Inches(0), h=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(13.333), h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def _text_box(slide, left, top, width, height, text, size=18, bold=False,
              color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _bullet_slide(slide, items, left=Inches(0.8), top=Inches(2.0),
                   width=Inches(11.5), size=20, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0


def _title_bar(slide, title, subtitle=None):
    _add_accent_bar(slide, top=Inches(0), h=Inches(1.4))
    _text_box(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
              title, size=32, bold=True, color=WHITE)
    if subtitle:
        _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
                  subtitle, size=16, color=RGBColor(0xCC, 0xDD, 0xFF))


def _add_image_safe(slide, img_name, left, top, width=None, height=None):
    path = FIG / img_name
    if not path.exists():
        for ext in [".png", ".pdf", ".jpg"]:
            alt = FIG / (Path(img_name).stem + ext)
            if alt.exists():
                path = alt
                break
    if path.exists() and path.suffix.lower() != ".pdf":
        kw = {}
        if width:
            kw["width"] = width
        if height:
            kw["height"] = height
        slide.shapes.add_picture(str(path), left, top, **kw)


# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── SLIDE 1: Title ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_add_bg(sl)
_add_accent_bar(sl, top=Inches(2.8), h=Inches(0.06))
_add_accent_bar(sl, top=Inches(4.6), h=Inches(0.04))

# Logo (real institute logo)
_add_image_safe(sl, "iiit_logo.png", Inches(5.4), Inches(0.3),
                width=Inches(2.5))

_text_box(sl, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.0),
          "DST-BA: Dynamic Service Time with Buffer-Aware\n"
          "UAV Data Collection for IoT Sensor Networks",
          size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

_text_box(sl, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.5),
          "B.Tech Project — Mid-Term Progress Report\n"
          "Rajneesh Chaudhary  |  2022IMG-052\n"
          "Advisor: Dr. Ankur Jaiswal\n"
          "ABV-IIITM Gwalior  •  March 2026",
          size=18, color=RGBColor(0xBB, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

# ── SLIDE 2: Agenda ────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Agenda")
items = [
    "1.  Problem Statement",
    "2.  Research Gaps",
    "3.  Proposed Framework (DST-BA)",
    "4.  Key Algorithms",
    "5.  Experimental Results",
    "6.  Analysis & Findings",
    "7.  Future Work",
]
_bullet_slide(sl, items, top=Inches(1.8), size=24)

# ── SLIDE 3: Problem Statement ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Problem Statement", "Why UAV-based IoT data collection?")
items = [
    "• IoT sensors in remote areas generate time-sensitive data",
    "• Multi-hop relay is costly — UAV provides efficient alternative",
    "• UAV flies close to sensors: better channel, faster data collection",
    "",
    "Objectives:",
    "  ✓  Visit ALL sensors with minimum energy",
    "  ✓  Realistic channel-aware communication",
    "  ✓  3D obstacle-safe trajectory planning",
]
_bullet_slide(sl, items, top=Inches(1.8), size=22)

# ── SLIDE 4: Research Gaps ──────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Research Gaps", "What existing work is missing")
items = [
    "1. Fixed hover time assumed — ignores real fading channels",
    "2. No RP compression — orphan nodes get missed",
    "3. Location-only clustering — ignores AoI & priority",
    "4. 2D-only paths — no 3D obstacle avoidance",
    "5. Ideal sensing assumed — no probabilistic model",
]
_bullet_slide(sl, items, top=Inches(1.8), size=22)

# ── SLIDE 5: Framework Overview ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "DST-BA: 8-Stage Pipeline")
items = [
    "1. Rendezvous Point Selection (reduce waypoints ~70%)",
    "2. Semantic Clustering (group by location + AoI + priority)",
    "3. GA Sequence Optimization (route ordering)",
    "4. PCA-GLS Path Refinement (smooth path)",
    "5. SCA Hover Optimization (best 3D position)",
    "6. DST-BA Service Time (channel-aware hover)",
    "7. TDMA Scheduling (interference-free)",
    "8. IoT Energy Tracking (monitor sensor budgets)",
]
_bullet_slide(sl, items, top=Inches(1.6), size=20)
_add_image_safe(sl, "routing_pipeline.png", Inches(8.5), Inches(1.8),
                width=Inches(4.5))

# ── SLIDE 6: Mathematical Models ───────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Key Mathematical Models")
items = [
    "UAV Energy:  E = E_propulsion + E_hover + E_communication",
    "",
    "Rician Channel:  Rate depends on distance, fading (K=15 dB)",
    "",
    "Service Time:  τ = max(sensing time, data/rate)",
    "               → adapts to channel conditions",
    "",
    "AoI Urgency:  older data gets higher collection priority",
]
_bullet_slide(sl, items, top=Inches(1.8), size=22)

# ── SLIDE 7: Design Choices ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Design Choices", "Why these algorithms?")
items = [
    "✓ DBSCAN Clustering — no need to pre-set # of clusters",
    "✓ Genetic Algorithm — handles NP-hard routing efficiently",
    "✓ Greedy Dominating Set — fast RP selection with coverage guarantee",
    "✓ Rician Fading (K=15 dB) — standard UAV air-to-ground model",
    "✓ SCA Hover — finds energy-optimal 3D hover positions",
]
_bullet_slide(sl, items, top=Inches(1.8), size=22)

# ── SLIDE 8: Implementation ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Implementation", "Python 3.13 — modular simulation platform")
items = [
    "• Modular architecture: config, core, path, metrics, visualization",
    "• Presets: 20 nodes (simple) and 50 nodes (full)",
    "• Single run, batch mode, comparison, ablation & scalability",
    "• Real-time 2D/3D visualization with animation export",
    "• Automated metric logging and LaTeX report generation",
]
_bullet_slide(sl, items, top=Inches(1.8), size=22)

# ── SLIDE 9: Environment & Trajectory ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl, RGBColor(0x10, 0x10, 0x18))
_title_bar(sl, "Environment & Trajectory")
_add_image_safe(sl, "environment.png", Inches(0.3), Inches(1.6),
                width=Inches(6.2))
_add_image_safe(sl, "trajectory_summary.png", Inches(6.8), Inches(1.6),
                width=Inches(6.2))

# ── SLIDE 10: 3D Trajectory ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl, RGBColor(0x10, 0x10, 0x18))
_title_bar(sl, "3D Trajectory Views")
_add_image_safe(sl, "trajectory_3d_isometric.png", Inches(0.3), Inches(1.6),
                width=Inches(6.2))
_add_image_safe(sl, "trajectory_3d_side_view.png", Inches(6.8), Inches(1.6),
                width=Inches(6.2))

# ── SLIDE 11: Algorithm Comparison ─────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Algorithm Comparison", "DST-BA vs. 4 baselines (seed=42)")

# Comparison table
from pptx.util import Inches, Pt
tbl_shape = sl.shapes.add_table(6, 7, Inches(0.4), Inches(1.6),
                                 Inches(12.5), Inches(2.8))
tbl = tbl_shape.table

headers = ["Algorithm", "Visited", "Coverage", "Energy (kJ)", "E/Node", "Steps", "Data (Mb)"]
data = [
    ["DST-BA (Proposed)", "19/19", "100.0%", "172.6", "9.1 kJ", "441", "285.4"],
    ["Nearest-Neighbour", "19/19", "100.0%", "352.0", "18.5 kJ", "253", "950.0"],
    ["Random-Walk", "18/19", "94.7%", "404.1", "22.4 kJ", "800", "900.0"],
    ["Fixed-Sweep", "19/19", "100.0%", "372.3", "19.6 kJ", "414", "950.0"],
    ["Cluster-First", "19/19", "100.0%", "352.0", "18.5 kJ", "253", "950.0"],
]

for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

for i, row in enumerate(data):
    for j, val in enumerate(row):
        cell = tbl.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE if i == 0 else DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        if i == 0:
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x50, 0x88)
        else:
            cell.fill.fore_color.rgb = LIGHT_GRAY if (i % 2) else WHITE

_add_image_safe(sl, "algorithm_comparison.png", Inches(0.4), Inches(4.6),
                width=Inches(12))

# ── SLIDE 12: Energy Analysis ──────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Energy Efficiency Analysis", "51–57% energy savings over all baselines")
_add_image_safe(sl, "energy_visited.png", Inches(0.3), Inches(1.6),
                width=Inches(6.2))
_add_image_safe(sl, "compare_energy-consumed-J.png", Inches(6.8), Inches(1.6),
                width=Inches(6.2))

# ── SLIDE 13: Coverage & Steps ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Coverage & Mission Duration")
_add_image_safe(sl, "compare_coverage-ratio-percent.png", Inches(0.3), Inches(1.6),
                width=Inches(6.2))
_add_image_safe(sl, "compare_steps.png", Inches(6.8), Inches(1.6),
                width=Inches(6.2))

# ── SLIDE 14: Stability & Communication ────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Stability & Communication Quality")
_add_image_safe(sl, "battery_replans.png", Inches(0.3), Inches(1.6),
                width=Inches(6.2))
_add_image_safe(sl, "communication_quality.png", Inches(6.8), Inches(1.6),
                width=Inches(6.2))

# ── SLIDE 15: Clustering & AoI ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Semantic Clustering & Age of Information")
_add_image_safe(sl, "semantic_clustering_geo.png", Inches(0.3), Inches(1.6),
                width=Inches(6.2))
_add_image_safe(sl, "aoi_timeline.png", Inches(6.8), Inches(1.6),
                width=Inches(6.2))

# ── SLIDE 16: Ablation Study ───────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Ablation Study", "What happens when we remove each component?")
items = [
    "Removed one component at a time:",
    "",
    "  • Without RP Compression → energy increases",
    "  • Without GA Routing → poor path quality",
    "  • Without SCA Hover → suboptimal positions",
]
_bullet_slide(sl, items, top=Inches(1.6), size=20,
              color=WHITE)
_add_image_safe(sl, "ablation_delta_bar.png", Inches(6.5), Inches(1.8),
                width=Inches(6.3))

# ── SLIDE 17: Scalability ──────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Scalability Analysis", "10 → 100 nodes")
_add_image_safe(sl, "scalability_node_count_replan_frequency.png",
                Inches(0.3), Inches(1.6), width=Inches(6.2))
_add_image_safe(sl, "scalability_node_count_path_stability_index.png",
                Inches(6.8), Inches(1.6), width=Inches(6.2))

# ── SLIDE 18: Dashboard & Radar ────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl, RGBColor(0x10, 0x10, 0x18))
_title_bar(sl, "Comprehensive Dashboard & Metrics")
_add_image_safe(sl, "dashboard_panel.png", Inches(0.3), Inches(1.6),
                width=Inches(6.2))
_add_image_safe(sl, "radar_chart.png", Inches(6.8), Inches(1.6),
                width=Inches(6.2))

# ── SLIDE 19: Key Findings ─────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Key Findings")
items = [
    "✓  100% sensor coverage (19/19 nodes visited)",
    "✓  172.6 kJ energy — 51–57% less than all baselines",
    "✓  71.2% battery remaining after full mission",
    "✓  Zero collisions, perfect path stability",
    "",
    "DST-BA halves energy consumption",
    "without sacrificing any coverage.",
]
_bullet_slide(sl, items, top=Inches(1.8), size=24)

# ── SLIDE 20: Limitations ──────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Current Limitations")
items = [
    "1. Single UAV only — multi-UAV not yet implemented",
    "2. Tested on 20-node scenario — scalability to 500+ pending",
    "3. No deep RL comparison yet (planned for next phase)",
]
_bullet_slide(sl, items, top=Inches(2.0), size=24)

# ── SLIDE 21: Future Work ──────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_title_bar(sl, "Future Work")
items = [
    "• Deep RL agent for real-time trajectory adaptation",
    "• Multi-UAV coordination for large-scale networks",
    "• Hardware validation with DJI drone + Raspberry Pi sensors",
    "• Advanced channels: NOMA, RIS, RF energy harvesting",
    "• Scale to 100–500 nodes with hierarchical clustering",
]
_bullet_slide(sl, items, top=Inches(1.8), size=22)

# ── SLIDE 22: Mission Progress ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl, RGBColor(0x10, 0x10, 0x18))
_title_bar(sl, "Mission Progress Overview")
_add_image_safe(sl, "mission_progress_combined.png", Inches(1.5), Inches(1.5),
                width=Inches(10))

# ── SLIDE 23: Thank You ────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_add_accent_bar(sl, top=Inches(2.8), h=Inches(0.06))
_add_accent_bar(sl, top=Inches(5.0), h=Inches(0.04))

_text_box(sl, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.8),
          "Thank You!", size=44, bold=True, color=WHITE,
          alignment=PP_ALIGN.CENTER)

_text_box(sl, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.6),
          "Questions & Discussion", size=28, color=ACCENT,
          alignment=PP_ALIGN.CENTER)

_text_box(sl, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.0),
          "Rajneesh Chaudhary  •  2022IMG-052\n"
          "Advisor: Dr. Ankur Jaiswal  •  ABV-IIITM Gwalior",
          size=18, color=RGBColor(0xAA, 0xBB, 0xDD),
          alignment=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Presentation saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
